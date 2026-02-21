"""
template_handler.py - YANGILANGAN VERSIYA
Yangi struktura: 1 sarlavha + 1 reja + n kontent + 1 rahmat
"""

import os
import glob
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import logging

logger = logging.getLogger(__name__)

TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")

TEMPLATE_FILENAMES = {
    "1_biznes_moliya": {i: f"{i}_*" for i in range(1, 11)},
    "2_talim_pedagogika": {i: f"{i}_*" for i in range(1, 11)},
    "3_texnologiya_it": {i: f"{i}_*" for i in range(1, 12)},
    "4_tibbiyot_soglik": {i: f"{i}_*" for i in range(1, 11)},
    "5_adabiyot_tilshunoslik": {i: f"{i}_*" for i in range(1, 11)},
    "6_ijtimoiy_fanlar": {i: f"{i}_*" for i in range(1, 12)},
    "7_fan_tadqiqot": {i: f"{i}_*" for i in range(1, 10)},
    "8_tarix_madaniyat": {i: f"{i}_*" for i in range(1, 11)},
    "9_tabiat_ekologiya": {i: f"{i}_*" for i in range(1, 11)},
    "10_kreativ_universal": {i: f"{i}_*" for i in range(1, 11)}
}


def find_template_file(category: str, template_id: int) -> str | None:
    """Template faylni topish"""
    try:
        category_dir = os.path.join(TEMPLATES_DIR, category.replace("_", " "))
        
        if not os.path.exists(category_dir):
            logger.error(f"Template papka topilmadi: {category_dir}")
            return None
        
        pattern = TEMPLATE_FILENAMES.get(category, {}).get(template_id)
        if not pattern:
            logger.error(f"Template ID topilmadi: {category}/{template_id}")
            return None
        
        search_pattern = os.path.join(category_dir, f"{pattern}.pptx")
        files = glob.glob(search_pattern)
        
        if files:
            logger.info(f"✅ Template topildi: {files[0]}")
            return files[0]
        else:
            logger.error(f"Template fayl topilmadi: {search_pattern}")
            return None
            
    except Exception as e:
        logger.error(f"Template qidirishda xato: {e}")
        return None


def fill_template_slides_new(template_path: str, content_data: dict, output_path: str) -> bool:
    """
    Template faylni AI content bilan to'ldirish - YANGILANGAN
    
    STRUKTURA:
    - SAHIFA 1: Sarlavha (title + author)
    - SAHIFA 2: Reja (title + plans)
    - SAHIFA 3 to n+2: Asosiy kontent
    - SAHIFA n+3: Rahmat (author)
    """
    try:
        prs = Presentation(template_path)
        logger.info(f"✅ Template ochildi: {len(prs.slides)} ta slayd")
        
        title = content_data.get("title", "")
        author = content_data.get("author", "")
        subtitle = content_data.get("subtitle", "")
        plans = content_data.get("plan", [])
        slides_data = content_data.get("slides", [])
        
        # SAHIFA 1: Sarlavha
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            _fill_title_slide(slide, title, author)
        
        # SAHIFA 2: Reja
        if len(prs.slides) > 1:
            slide = prs.slides[1]
            _fill_plan_slide(slide, title, plans)
        
        # SAHIFA 3 to n+2: Kontent
        for idx, slide_content in enumerate(slides_data):
            slide_idx = idx + 2  # +2 chunki 0=title, 1=plan
            if slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                _fill_content_slide(slide, slide_content)
        
        # SAHIFA n+3: Rahmat
        last_idx = len(slides_data) + 2
        if last_idx < len(prs.slides):
            slide = prs.slides[last_idx]
            _fill_thanks_slide(slide, author)
        
        prs.save(output_path)
        logger.info(f"✅ Taqdimot saqlandi: {output_path}")
        return True
        
    except Exception as e:
        logger.error(f"Template to'ldirishda xato: {e}")
        return False


def _fill_title_slide(slide, title: str, author: str):
    """SAHIFA 1: Sarlavha + Author"""
    try:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                
                # Title
                if ph_type in [1, 3]:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = title
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(44)
                
                # Subtitle/Author
                elif ph_type == 2:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = author
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.size = Pt(20)
                        
    except Exception as e:
        logger.warning(f"Title slide xatosi: {e}")


def _fill_plan_slide(slide, title: str, plans: list):
    """SAHIFA 2: Reja"""
    try:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                
                # Title
                if ph_type == 1:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = title
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(32)
                
                # Content - Rejalar
                elif ph_type == 2:
                    text_frame.clear()
                    for i, plan in enumerate(plans[:3], 1):
                        p = text_frame.add_paragraph()
                        p.text = f"{i}. {plan}"
                        p.level = 0
                        p.space_before = Pt(6)
                        for run in p.runs:
                            run.font.size = Pt(20)
                            run.font.bold = True
                            
    except Exception as e:
        logger.warning(f"Plan slide xatosi: {e}")


def _fill_content_slide(slide, content: dict):
    """SAHIFA 3-n+2: Kontent"""
    try:
        title = content.get("title", "")
        points = content.get("points", [])
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                
                # Title
                if ph_type == 1:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = title
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(32)
                
                # Content
                elif ph_type == 2:
                    text_frame.clear()
                    for point in points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.space_before = Pt(6)
                        for run in p.runs:
                            run.font.size = Pt(18)
                            
    except Exception as e:
        logger.warning(f"Content slide xatosi: {e}")


def _fill_thanks_slide(slide, author: str):
    """SAHIFA n+3: Rahmat"""
    try:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                
                # Title
                if ph_type in [1, 3]:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = "Etiboringiz uchun rahmat!"
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(40)
                
                # Subtitle/Author
                elif ph_type == 2:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = author
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.size = Pt(24)
                        
    except Exception as e:
        logger.warning(f"Thanks slide xatosi: {e}")


def create_presentation_from_template(
    category: str,
    template_id: int,
    content_data: dict,
    output_filename: str
) -> str | None:
    """Template asosida taqdimot yaratish"""
    try:
        template_path = find_template_file(category, template_id)
        if not template_path:
            logger.error("Template fayl topilmadi!")
            return None
        
        output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, output_filename)
        
        success = fill_template_slides_new(template_path, content_data, output_path)
        
        if success:
            return output_path
        else:
            return None
            
    except Exception as e:
        logger.error(f"Presentation yaratishda xato: {e}")
        return None

print("✅ Yangi template_handler tayyor")
