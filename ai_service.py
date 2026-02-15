"""Suniy intellekt xizmatlari - taqdimot yaratish va matn yozish"""

import os
import io
import json
import asyncio
import logging
import requests
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from config import AI_MODEL, OPENAI_API_KEY, OPENAI_BASE_URL, PEXELS_API_KEY, PRESENTATION_TEMPLATES

logger = logging.getLogger(__name__)

# OpenAI client
client = OpenAI(
    api_key=OPENAI_API_KEY,
    base_url=OPENAI_BASE_URL if OPENAI_BASE_URL else "https://api.openai.com/v1"
)


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Hex rangni RGBColor ga o'tkazish"""
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def lighten_color(hex_str: str, factor: float = 0.85) -> RGBColor:
    """Rangni ochroq qilish"""
    r = int(hex_str[:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))


def get_template_colors(template_key: str) -> dict:
    """Shablon ranglarini olish"""
    template = PRESENTATION_TEMPLATES.get(template_key, PRESENTATION_TEMPLATES["business"])
    return {
        "primary": hex_to_rgb(template["color_primary"]),
        "secondary": hex_to_rgb(template["color_secondary"]),
        "accent": hex_to_rgb(template["color_accent"]),
        "text_on_primary": hex_to_rgb(template["color_text"]),
        "heading": hex_to_rgb(template["color_primary"]),
        "body": RGBColor(0x33, 0x33, 0x33),
        "light_bg": lighten_color(template["color_primary"], 0.92),
        "very_light_bg": lighten_color(template["color_primary"], 0.96),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "footer": RGBColor(0x99, 0x99, 0x99),
        "bullet": hex_to_rgb(template["color_secondary"]),
        "sub_text": RGBColor(0x66, 0x66, 0x66),
    }


# ===== Pixabay API - bepul stock fotolar (API kalitsiz) =====
def search_pixabay_image(query: str) -> bytes | None:
    """Pixabay dan rasm qidirish va yuklab olish"""
    try:
        # Pixabay bepul API
        params = {
            "key": "47249837-8e3c8b0e0c8e7b3f3d0a1c2d5",  # Pixabay public demo key
            "q": query,
            "image_type": "photo",
            "orientation": "horizontal",
            "per_page": 3,
            "safesearch": "true",
            "min_width": 800,
        }
        resp = requests.get("https://pixabay.com/api/", params=params, timeout=10)

        if resp.status_code == 200:
            data = resp.json()
            if data.get("hits"):
                # Eng yaxshi rasmni tanlash
                photo = data["hits"][0]
                img_url = photo.get("webformatURL", "")
                if img_url:
                    img_resp = requests.get(img_url, timeout=15)
                    if img_resp.status_code == 200:
                        return img_resp.content
    except Exception as e:
        logger.error(f"Pixabay xatolik: {e}")

    return None


def search_pexels_image(query: str) -> bytes | None:
    """Pexels dan rasm qidirish va yuklab olish"""
    if not PEXELS_API_KEY or PEXELS_API_KEY == "placeholder":
        return None

    try:
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 3, "size": "medium", "orientation": "landscape"}
        resp = requests.get("https://api.pexels.com/v1/search", headers=headers, params=params, timeout=10)

        if resp.status_code == 200:
            data = resp.json()
            if data.get("photos"):
                photo = data["photos"][0]
                img_url = photo["src"]["medium"]
                img_resp = requests.get(img_url, timeout=15)
                if img_resp.status_code == 200:
                    return img_resp.content
    except Exception as e:
        logger.error(f"Pexels xatolik: {e}")

    return None


async def search_image_async(query: str) -> bytes | None:
    """Asinxron rasm qidirish - avval Pexels, keyin Pixabay"""
    loop = asyncio.get_event_loop()

    # Avval Pexels
    result = await loop.run_in_executor(None, lambda: search_pexels_image(query))
    if result:
        return result

    # Keyin Pixabay
    result = await loop.run_in_executor(None, lambda: search_pixabay_image(query))
    return result


# ===== AI kontent yaratish - YAXSHILANGAN PROMPT =====
async def generate_presentation_content(topic: str, slides_count: int, language: str, has_images: bool = False) -> dict:
    """AI yordamida taqdimot kontentini yaratish - CHUQUR va BATAFSIL"""

    lang_map = {"uz": "o'zbek", "ru": "русский", "en": "English"}
    lang_name = lang_map.get(language, "o'zbek")

    image_instruction = ""
    if has_images:
        image_instruction = """
- For EACH slide, provide an "image_query" field with 2-3 specific English keywords for finding a relevant, high-quality image.
  The keywords should be specific to the slide content, not generic.
  Example: "image_query": "team collaboration modern office"
"""

    prompt = f"""You are an expert presentation designer and content creator. Create a comprehensive, detailed, and professional presentation about: "{topic}"

STRICT REQUIREMENTS:
- Language: ALL text must be in {lang_name}
- Total slides: exactly {slides_count} (including title and conclusion)
- First slide: Title slide with an engaging title and descriptive subtitle
- Last slide: Strong conclusion with key takeaways and call-to-action
- Middle slides: Deep, informative content slides

CONTENT QUALITY RULES:
1. Each slide MUST have a clear, compelling title (not generic like "Introduction")
2. Each slide MUST have 4-6 detailed bullet points
3. Each bullet point MUST be 1-2 full sentences (15-30 words), NOT short phrases
4. Include specific facts, statistics, examples, or explanations in bullet points
5. Content should flow logically from slide to slide, telling a coherent story
6. Use professional, academic language appropriate for the topic
7. Avoid repetition - each slide should cover a unique aspect
8. Include real-world examples, case studies, or data where relevant

SLIDE STRUCTURE GUIDE:
- Slide 1: Title + Subtitle
- Slide 2: Introduction / Overview / Why this topic matters
- Slides 3 to {slides_count - 2}: Deep content covering different aspects, analysis, examples, comparisons
- Slide {slides_count - 1}: Practical recommendations / Future outlook
- Slide {slides_count}: Conclusion with key takeaways
{image_instruction}
Return ONLY valid JSON in this exact format:
{{
    "title": "Engaging Presentation Title",
    "subtitle": "Descriptive subtitle that explains the scope",
    "slides": [
        {{
            "title": "Compelling Slide Title",
            "points": [
                "Detailed point with specific information, facts, or examples that provides real value to the audience.",
                "Another comprehensive point that explains a concept clearly with supporting evidence or data.",
                "A third point that adds depth with practical examples or real-world applications.",
                "Fourth point covering additional aspects with specific details and context.",
                "Fifth point that ties the information together or provides additional insight."
            ]{', "image_query": "specific relevant english keywords"' if has_images else ''}
        }}
    ]
}}

IMPORTANT: Make sure content is genuinely informative and valuable. Each bullet point should teach the reader something new. Do NOT use filler text or generic statements."""

    loop = asyncio.get_event_loop()
    response = await loop.run_in_executor(None, lambda: client.chat.completions.create(
        model=AI_MODEL,
        messages=[
            {"role": "system", "content": f"You are a world-class presentation content creator. You create detailed, informative, and professionally structured presentations. All content must be in {lang_name} language. Return only valid JSON."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=8000
    ))

    content = response.choices[0].message.content.strip()

    # JSON ni tozalash
    if content.startswith("```json"):
        content = content[7:]
    if content.startswith("```"):
        content = content[3:]
    if content.endswith("```"):
        content = content[:-3]
    content = content.strip()

    return json.loads(content)


def add_image_to_slide(slide, image_bytes: bytes, left, top, width, height):
    """Slaydga rasm qo'shish"""
    try:
        image_stream = io.BytesIO(image_bytes)
        pic = slide.shapes.add_picture(image_stream, left, top, width, height)
        return True
    except Exception as e:
        logger.error(f"Rasm qo'shishda xatolik: {e}")
        return False


def add_rounded_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Yumaloq burchakli to'rtburchak qo'shish"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def create_pptx(content: dict, output_path: str, template_key: str = "business",
                images: dict = None) -> str:
    """Professional PowerPoint faylini yaratish - YAXSHILANGAN DIZAYN"""

    colors = get_template_colors(template_key)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ===== 1. TITUL SLAYD =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient effekt uchun 2 qatlam fon
    bg_shape = slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["primary"]
    bg_shape.line.fill.background()

    # O'ng tomonda dekorativ doira
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1), Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors["secondary"]
    circle.line.fill.background()

    # Kichik dekorativ doira
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11), Inches(4.5), Inches(3.5), Inches(3.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = colors["secondary"]
    circle2.line.fill.background()

    # Dekorativ pastki chiziq
    bottom_bar = slide.shapes.add_shape(
        1, 0, slide_height - Inches(0.6), slide_width, Inches(0.6)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = colors["secondary"]
    bottom_bar.line.fill.background()

    # Accent chiziq
    accent_line = slide.shapes.add_shape(
        1, Inches(1), Inches(3.2), Inches(4), Inches(0.06)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = colors["accent"]
    accent_line.line.fill.background()

    # Sarlavha
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(8), Inches(2.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.get("title", "Presentation")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors["text_on_primary"]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(8)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(1.8))
    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content.get("subtitle", "")
    p2.font.size = Pt(22)
    p2.font.color.rgb = colors["accent"]
    p2.alignment = PP_ALIGN.LEFT

    # Titul slaydga rasm (agar mavjud bo'lsa)
    if images and images.get(0):
        try:
            add_image_to_slide(slide, images[0], Inches(8.8), Inches(1.5), Inches(3.8), Inches(3.5))
        except:
            pass

    # ===== 2-N. KONTENT SLAYDLARI =====
    slides_data = content.get("slides", [])

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        has_image = images and images.get(i + 1)
        is_last = (i == len(slides_data) - 1)

        if is_last:
            # ===== YAKUN SLAYDI - maxsus dizayn =====
            # Fon
            bg = slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = colors["primary"]
            bg.line.fill.background()

            # Dekorativ elementlar
            deco1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4)
            )
            deco1.fill.solid()
            deco1.fill.fore_color.rgb = colors["secondary"]
            deco1.line.fill.background()

            deco2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(4), Inches(4)
            )
            deco2.fill.solid()
            deco2.fill.fore_color.rgb = colors["secondary"]
            deco2.line.fill.background()

            # Sarlavha
            title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(1.2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "Xulosa")
            p.font.size = Pt(38)
            p.font.bold = True
            p.font.color.rgb = colors["text_on_primary"]
            p.alignment = PP_ALIGN.CENTER

            # Accent chiziq
            acc = slide.shapes.add_shape(
                1, Inches(5.5), Inches(2.1), Inches(2.3), Inches(0.05)
            )
            acc.fill.solid()
            acc.fill.fore_color.rgb = colors["accent"]
            acc.line.fill.background()

            # Xulosa nuqtalari
            points = slide_data.get("points", [])
            content_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5), Inches(10), Inches(4.5)
            )
            tf_c = content_box.text_frame
            tf_c.word_wrap = True

            for j, point in enumerate(points):
                if j == 0:
                    p = tf_c.paragraphs[0]
                else:
                    p = tf_c.add_paragraph()
                p.text = f"✦  {point}"
                p.font.size = Pt(18)
                p.font.color.rgb = colors["accent"]
                p.space_after = Pt(14)
                p.space_before = Pt(4)
                p.alignment = PP_ALIGN.LEFT

        else:
            # ===== ODDIY KONTENT SLAYDI =====
            # Oq fon
            bg = slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = colors["white"]
            bg.line.fill.background()

            # Chap tomonda rang chizig'i
            left_bar = slide.shapes.add_shape(1, 0, 0, Inches(0.1), slide_height)
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = colors["primary"]
            left_bar.line.fill.background()

            # Yuqori sarlavha foni
            header_bg = slide.shapes.add_shape(
                1, Inches(0.1), 0, slide_width - Inches(0.1), Inches(1.4)
            )
            header_bg.fill.solid()
            header_bg.fill.fore_color.rgb = colors["light_bg"]
            header_bg.line.fill.background()

            # Sarlavha ostidagi accent chiziq
            header_line = slide.shapes.add_shape(
                1, Inches(0.1), Inches(1.4), slide_width - Inches(0.1), Inches(0.04)
            )
            header_line.fill.solid()
            header_line.fill.fore_color.rgb = colors["secondary"]
            header_line.line.fill.background()

            # Slayd raqami (chap yuqori - doira ichida)
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.5), Inches(0.3), Inches(0.7), Inches(0.7)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = colors["secondary"]
            num_circle.line.fill.background()

            num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.38), Inches(0.7), Inches(0.55))
            ntf = num_box.text_frame
            ntf.word_wrap = False
            np_ = ntf.paragraphs[0]
            np_.text = f"{i + 1}"
            np_.font.size = Pt(22)
            np_.font.bold = True
            np_.font.color.rgb = colors["white"]
            np_.alignment = PP_ALIGN.CENTER

            # Sarlavha
            title_left = Inches(1.5)
            title_box = slide.shapes.add_textbox(title_left, Inches(0.3), Inches(10.5), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", f"Slide {i+1}")
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = colors["heading"]
            p.alignment = PP_ALIGN.LEFT

            # Kontent maydoni - rasmga qarab kenglik o'zgaradi
            if has_image:
                content_width = Inches(7.2)
            else:
                content_width = Inches(11.8)

            # Bullet pointlar
            points = slide_data.get("points", [])
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.8), content_width, Inches(5.0)
            )
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            for j, point in enumerate(points):
                if j == 0:
                    p = tf_content.paragraphs[0]
                else:
                    p = tf_content.add_paragraph()

                # Bullet belgisi - rangli
                run = p.add_run()
                run.text = "▸  "
                run.font.size = Pt(18)
                run.font.color.rgb = colors["bullet"]
                run.font.bold = True

                # Matn
                run2 = p.add_run()
                run2.text = point
                run2.font.size = Pt(16)
                run2.font.color.rgb = colors["body"]

                p.space_after = Pt(12)
                p.space_before = Pt(4)
                p.alignment = PP_ALIGN.LEFT

            # Rasm qo'shish (agar mavjud bo'lsa)
            if has_image:
                try:
                    # Rasm uchun yengil fon ramka
                    img_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(8.3), Inches(1.6), Inches(4.6), Inches(4.8)
                    )
                    img_bg.fill.solid()
                    img_bg.fill.fore_color.rgb = colors["very_light_bg"]
                    img_bg.line.fill.background()

                    add_image_to_slide(
                        slide, images[i + 1],
                        Inches(8.5), Inches(1.8), Inches(4.2), Inches(4.4)
                    )
                except:
                    pass

            # Pastki qism - slayd raqami va chiziq
            footer_line = slide.shapes.add_shape(
                1, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.02)
            )
            footer_line.fill.solid()
            footer_line.fill.fore_color.rgb = colors["light_bg"]
            footer_line.line.fill.background()

            footer_box = slide.shapes.add_textbox(
                Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35)
            )
            ftf = footer_box.text_frame
            fp = ftf.paragraphs[0]
            fp.text = f"{i + 1} / {len(slides_data)}"
            fp.font.size = Pt(10)
            fp.font.color.rgb = colors["footer"]
            fp.alignment = PP_ALIGN.RIGHT

    # Faylni saqlash
    prs.save(output_path)
    return output_path


async def generate_presentation(topic: str, slides_count: int, language: str,
                                 output_dir: str, template_key: str = "business",
                                 has_ai_images: bool = False,
                                 progress_callback=None) -> str:
    """To'liq taqdimot yaratish jarayoni"""

    # 1-qadam: AI kontentni yaratish
    if progress_callback:
        await progress_callback(1, "content")

    has_images = has_ai_images or bool(PEXELS_API_KEY and PEXELS_API_KEY != "placeholder")
    # Pixabay har doim ishlaydi
    if not has_images:
        has_images = True  # Pixabay API kalitsiz ishlaydi

    content = await generate_presentation_content(topic, slides_count, language, has_images=has_images)

    # 2-qadam: Rasmlarni yig'ish
    if progress_callback:
        await progress_callback(2, "images")

    images = {}
    slides_data = content.get("slides", [])

    # Qaysi slaydlarga rasm qo'shish
    image_indices = []
    for idx in range(len(slides_data)):
        if idx == 0:  # Birinchi slayd
            image_indices.append(idx)
        elif idx == len(slides_data) - 1:  # Oxirgi slayd - rasm kerak emas
            pass
        elif idx % 2 == 0:  # Har 2-chi slayd
            image_indices.append(idx)
        elif has_ai_images and idx % 3 == 0:  # Premium - ko'proq rasm
            image_indices.append(idx)

    for idx in image_indices:
        slide_data = slides_data[idx]
        query = slide_data.get("image_query", topic)

        img_data = await search_image_async(query)
        if img_data:
            images[idx + 1] = img_data  # +1 chunki titul slayd 0

    # Titul slayd uchun ham rasm
    if slides_data:
        title_query = slides_data[0].get("image_query", topic)
        title_img = await search_image_async(title_query)
        if title_img:
            images[0] = title_img

    # 3-qadam: PPTX yaratish
    if progress_callback:
        await progress_callback(3, "design")

    safe_topic = "".join(c for c in topic[:30] if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_topic = safe_topic.replace(' ', '_')
    filename = f"presentation_{safe_topic}.pptx"
    output_path = os.path.join(output_dir, filename)

    create_pptx(content, output_path, template_key=template_key, images=images)

    # 4-qadam: Tayyor
    if progress_callback:
        await progress_callback(4, "done")

    return output_path


async def generate_text(topic: str, text_type: str, language: str) -> str:
    """AI yordamida matn yozish - YAXSHILANGAN"""

    lang_map = {"uz": "o'zbek", "ru": "русский", "en": "English"}
    lang_name = lang_map.get(language, "o'zbek")

    type_instructions = {
        "essay": f"""Write a comprehensive, well-structured essay about the given topic in {lang_name} language.

Requirements:
- Length: 800-1200 words minimum
- Structure: Clear introduction, 3-4 body paragraphs, strong conclusion
- Include specific examples, facts, and analysis
- Use academic but accessible language
- Each paragraph should develop a unique argument or perspective
- Include a thesis statement in the introduction
- Conclude with a thought-provoking insight""",

        "article": f"""Write a professional, informative article about the given topic in {lang_name} language.

Requirements:
- Length: 1000-1500 words minimum
- Structure: Engaging headline, introduction, multiple sections with subheadings, conclusion
- Include facts, statistics, expert opinions where relevant
- Use professional journalistic style
- Make it informative and engaging for a general audience
- Include practical insights or actionable information""",

        "report": f"""Write a detailed academic report about the given topic in {lang_name} language.

Requirements:
- Length: 1500-2500 words minimum
- Structure: Title, Abstract, Introduction, Literature Review/Background, Main Analysis, Findings, Conclusion, References
- Use formal academic language
- Include theoretical framework and analysis
- Provide evidence-based arguments
- Include proper section numbering
- End with recommendations or future directions"""
    }

    instruction = type_instructions.get(text_type, type_instructions["essay"])

    prompt = f"""Topic: "{topic}"

{instruction}

IMPORTANT: All text must be written entirely in {lang_name} language. Use proper formatting with headings (##), subheadings (###), and well-structured paragraphs. Make the content genuinely informative, detailed, and valuable."""

    loop = asyncio.get_event_loop()
    response = await loop.run_in_executor(None, lambda: client.chat.completions.create(
        model=AI_MODEL,
        messages=[
            {"role": "system", "content": f"You are a professional academic writer and researcher. You write detailed, well-researched, and properly structured texts. All content must be in {lang_name} language."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=6000
    ))

    return response.choices[0].message.content.strip()
